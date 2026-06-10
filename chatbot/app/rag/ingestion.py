"""Ingestion: scrape the portfolio site and load project markdown into Chroma."""

import ipaddress
import re
import socket
from pathlib import Path
from urllib.parse import urljoin, urlparse

import httpx
from bs4 import BeautifulSoup, NavigableString
from langchain_core.documents import Document
from langchain_text_splitters import (
    MarkdownHeaderTextSplitter,
    RecursiveCharacterTextSplitter,
)

from app.config import get_settings
from app.rag.vectorstore import add_documents_safe, get_vectorstore
from app.utils.sanitizer import sanitize_user_message

TABLE_PATTERN = re.compile(r"^\s*\|.+\|\s*$", re.M)
FORMULA_PATTERN = re.compile(r"(?ms)(\$\$.*?\$\$|`{3,}math.*?`{3,}|\$[^\$\n]+\$)")
IMAGE_PATTERN = re.compile(r"!\[([^\]]*)\]\(([^)\s]+)(?:\s+\"([^\"]*)\")?\)")
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".svg", ".webp"}


def _try_image_ocr_bytes(image_bytes: bytes) -> str:
    try:
        from io import BytesIO

        import pytesseract
        from PIL import Image
    except ImportError:
        return ""

    try:
        with Image.open(BytesIO(image_bytes)) as img:
            text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception:
        return ""


def _try_image_ocr(image_path: Path) -> str:
    try:
        return _try_image_ocr_bytes(image_path.read_bytes())
    except Exception:
        return ""


def _try_image_vision_caption_bytes(
    image_bytes: bytes, filename: str = "image.png"
) -> str:
    try:
        import base64

        import openai
    except ImportError:
        return ""

    try:
        settings = get_settings()
        encoded = base64.b64encode(image_bytes).decode("ascii")
        suffix = Path(filename).suffix.lstrip(".").lower() or "png"
        image_url = f"data:image/{suffix};base64,{encoded}"
        client = openai.OpenAI(api_key=settings.openai_api_key)
        response = client.responses.create(
            model=settings.openai_vision_model,
            input=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "input_text",
                            "text": (
                                "Describe this image in one concise sentence, "
                                "focusing on the architecture or diagram content."
                            ),
                        },
                        {
                            "type": "input_image",
                            "image_url": image_url,
                        },
                    ],
                }
            ],
            max_output_tokens=100,
            temperature=0.0,
        )
        return response.output_text.strip()
    except Exception:
        return ""


def _try_image_vision_caption(image_path: Path) -> str:
    try:
        return _try_image_vision_caption_bytes(
            image_path.read_bytes(), filename=image_path.name
        )
    except Exception:
        return ""


def _download_image_bytes(image_url: str, client: httpx.Client) -> bytes | None:
    try:
        resp = client.get(image_url, follow_redirects=True, timeout=30)
        resp.raise_for_status()
        content_type = resp.headers.get("content-type", "")
        if not content_type.startswith("image/"):
            return None
        return resp.content
    except Exception:
        return None


def _collect_web_image_docs(
    soup: BeautifulSoup,
    page_url: str,
    page_title: str,
    client: httpx.Client,
) -> list[Document]:
    docs: list[Document] = []
    for img in soup.find_all("img"):
        src = img.get("src", "").strip()
        if not src or src.startswith("data:"):
            continue

        image_url = urljoin(page_url, src)
        image_bytes = _download_image_bytes(image_url, client)
        if not image_bytes:
            continue

        vision_caption = _try_image_vision_caption_bytes(
            image_bytes, filename=image_url
        )
        alt = img.get("alt", "").strip()
        title = img.get("title", "").strip()
        filename = (
            Path(urlparse(image_url).path).stem.replace("_", " ").replace("-", " ")
        )
        caption = vision_caption or alt or title or filename
        caption_label = "Vision caption" if vision_caption else "Caption"
        ocr_text = _try_image_ocr_bytes(image_bytes)
        content = (
            f"Image file: {image_url}. "
            f"Page: {page_title}. "
            f"{caption_label}: {caption}."
        )
        if ocr_text:
            content += f" OCR text: {ocr_text}."

        doc = Document(
            page_content=content,
            metadata={
                "source": page_url,
                "image_url": image_url,
                "page_title": page_title,
                "type": "website_image",
                "has_image": True,
                "vision_caption": bool(vision_caption),
                "caption_source": "vision" if vision_caption else "html",
                "ocr_text": bool(ocr_text),
            },
        )
        docs.append(doc)
    return docs


def _extract_table_text(table) -> str:
    rows = []
    for tr in table.find_all("tr"):
        cells = [cell.get_text(" ", strip=True) for cell in tr.find_all(["th", "td"])]
        if cells:
            rows.append(" | ".join(cells))
    if not rows:
        return ""
    return "Table:\n" + "\n".join(rows)


def _replace_html_structures(soup: BeautifulSoup) -> None:
    for table in soup.find_all("table"):
        table_text = _extract_table_text(table)
        table.replace_with(NavigableString(table_text or ""))

    for figure in soup.find_all("figure"):
        caption = figure.find("figcaption")
        if caption and caption.get_text(strip=True):
            caption_text = f"Caption: {caption.get_text(' ', strip=True)}"
            figure.insert_after(NavigableString(caption_text))
        figure.unwrap()

    math_tags = soup.find_all(["math"])
    for tag in math_tags:
        formula_text = tag.get_text(" ", strip=True)
        tag.replace_with(NavigableString(f"Formula:\n{formula_text}"))

    for img in soup.find_all("img"):
        alt = img.get("alt", "").strip()
        title = img.get("title", "").strip()
        src = img.get("src", "").strip()
        image_text = f"Image: alt='{alt}' src='{src}'"
        if title:
            image_text += f" title='{title}'"
        img.replace_with(NavigableString(image_text))


def _collect_local_image_docs(
    text: str, project_dir: Path, project_name: str
) -> list[Document]:
    docs: list[Document] = []
    for alt, src, title in IMAGE_PATTERN.findall(text):
        if not src or src.startswith("http://") or src.startswith("https://"):
            continue
        image_path = Path(src)
        if not image_path.is_absolute():
            image_path = project_dir / image_path
        if not image_path.exists() or image_path.suffix.lower() not in IMAGE_EXTENSIONS:
            continue

        vision_caption = _try_image_vision_caption(image_path)
        caption = (
            vision_caption.strip()
            or alt.strip()
            or title.strip()
            or image_path.stem.replace("_", " ").replace("-", " ")
        )
        caption_label = "Vision caption" if vision_caption else "Caption"
        ocr_text = _try_image_ocr(image_path)
        content = (
            f"Image file: {image_path.name}. "
            f"Project: {project_name}. "
            f"{caption_label}: {caption}."
        )
        if ocr_text:
            content += f" OCR text: {ocr_text}."
        doc = Document(
            page_content=content,
            metadata={
                "source": str(image_path),
                "project": project_name,
                "type": "image",
                "has_image": True,
                "vision_caption": bool(vision_caption),
                "caption_source": "vision" if vision_caption else "markdown",
                "ocr_text": bool(ocr_text),
            },
        )
        docs.append(doc)
    return docs


def _normalize_markdown(text: str) -> str:
    def _replace_formula(match: re.Match) -> str:
        formula = match.group(0).strip()
        if formula.startswith("$$") and formula.endswith("$$"):
            body = formula[2:-2].strip()
        elif formula.startswith("`" * 3):
            body = formula.strip("`").replace("math", "", 1).strip()
        else:
            body = formula.strip("$")
        return f"\n\nFormula:\n{body}\n\n"

    def _replace_image(match: re.Match) -> str:
        alt, src, title = match.groups()
        caption = f"Image: alt='{alt or ''}' src='{src}'"
        if title:
            caption += f" title='{title}'"
        return f"\n\n{caption}\n\n"

    text = FORMULA_PATTERN.sub(_replace_formula, text)
    text = IMAGE_PATTERN.sub(_replace_image, text)
    return text


def _chunk_metadata(doc: Document) -> None:
    content = doc.page_content
    doc.metadata.setdefault("has_table", bool(TABLE_PATTERN.search(content)))
    doc.metadata.setdefault("has_formula", bool(FORMULA_PATTERN.search(content)))
    doc.metadata.setdefault("has_image", "Image:" in content)


# --------------------------------------------------------------------------- #
# Website
# --------------------------------------------------------------------------- #


def scrape_website(urls: list[str], timeout: int = 30) -> list[Document]:
    """Fetch each URL and extract clean visible text.

    Strips <script>, <style>, <nav>, <header>, <footer>. If your site is a
    SPA where content is rendered client-side, swap httpx for Playwright.
    """
    docs: list[Document] = []

    def _is_safe_url(u: str) -> bool:
        try:
            parsed = urlparse(u)
            if parsed.scheme not in ("http", "https"):
                return False
            host = parsed.hostname
            if not host:
                return False
            # If hostname is an IP literal, block private/loopback/link-local
            try:
                ip = ipaddress.ip_address(host)
                if ip.is_private or ip.is_loopback or ip.is_link_local:
                    return False
            except ValueError:
                # not an IP literal; attempt dns lookup and check resolved IP
                try:
                    addr = socket.gethostbyname(host)
                    ip = ipaddress.ip_address(addr)
                    if ip.is_private or ip.is_loopback or ip.is_link_local:
                        return False
                except Exception:
                    # If DNS fails, be conservative and allow (could be public domain)
                    pass
            return True
        except Exception:
            return False

    with httpx.Client(
        timeout=timeout,
        follow_redirects=True,
        headers={"User-Agent": "PortfolioBot/1.0"},
    ) as client:
        for url in urls:
            if not _is_safe_url(url):
                print(f"[scrape] skipping unsafe url: {url}")
                continue
            try:
                resp = client.get(url)
                resp.raise_for_status()
            except Exception as e:
                print(f"[scrape] {url} failed: {e}")
                continue

            soup = BeautifulSoup(resp.text, "lxml")
            for tag in soup(["script", "style", "nav", "header", "footer", "noscript"]):
                tag.decompose()

            _replace_html_structures(soup)
            title = (
                soup.title.string.strip() if soup.title and soup.title.string else url
            )
            image_docs = _collect_web_image_docs(soup, url, title, client)
            text = soup.get_text(separator="\n", strip=True)
            if not text and not image_docs:
                print(f"[scrape] {url} produced no text and no images — skipping")
                continue

            if text:
                doc = Document(
                    page_content=text,
                    metadata={
                        "source": url,
                        "title": title,
                        "domain": urlparse(url).netloc,
                        "type": "website",
                    },
                )
                _chunk_metadata(doc)
                docs.append(doc)
                print(f"[scrape] {url} -> {len(text)} chars")

            for image_doc in image_docs:
                _chunk_metadata(image_doc)
                docs.append(image_doc)
    return docs


# --------------------------------------------------------------------------- #
# Projects: File Extraction
# --------------------------------------------------------------------------- #


def _extract_docx_text(path: Path) -> str:
    """Extract text from a DOCX file, preserving document structure with markdown headers.

    Converts Word heading styles to markdown headers (Heading 1 -> #, Heading 2 -> ##, etc.)
    and tables to pipe-delimited format. This structure is later used by MarkdownHeaderTextSplitter
    to preserve document hierarchy during chunking.
    """
    try:
        from docx import Document as DocxDocument
    except ImportError:
        print(f"[docx] python-docx not installed; skipping {path.name}")
        return ""

    try:
        doc = DocxDocument(path)
        text_parts = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Map Word heading styles to markdown headers
            style_name = para.style.name if para.style else ""
            if "Heading 1" in style_name:
                text_parts.append(f"# {text}")
            elif "Heading 2" in style_name:
                text_parts.append(f"## {text}")
            elif "Heading 3" in style_name:
                text_parts.append(f"### {text}")
            elif "Heading 4" in style_name:
                text_parts.append(f"#### {text}")
            else:
                text_parts.append(text)

        # Extract table content
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                text_parts.append("Table:\n" + "\n".join(rows))

        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"[docx] Failed to extract from {path.name}: {e}")
        return ""


def _extract_pptx_text(path: Path) -> tuple[str, list[bytes]]:
    """Extract text and images from a PPTX file.

    Returns: (markdown_text_with_headers, list_of_image_bytes)
    Converts slide content to markdown with slide numbers as headers.
    """
    try:
        from pptx import Presentation
    except ImportError:
        print(f"[pptx] python-pptx not installed; skipping {path.name}")
        return "", []

    try:
        prs = Presentation(path)
        text_parts = []
        image_bytes_list = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            text_parts.append(f"# Slide {slide_num}")

            # Extract text from shapes (titles, text boxes)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

                # Extract embedded images
                if shape.shape_type == 13:  # PICTURE
                    try:
                        image = shape.image
                        image_bytes_list.append(image.blob)
                    except Exception:
                        pass

            text_parts.append("")  # Blank line between slides

        return "\n\n".join(text_parts), image_bytes_list
    except Exception as e:
        print(f"[pptx] Failed to extract from {path.name}: {e}")
        return "", []


def _load_standalone_images(
    projects_dir: Path, exclude_files: set[str]
) -> list[Document]:
    """Load all standalone .png files in projects_dir with vision captions and OCR.

    Excludes images already referenced in markdown files (passed in exclude_files).
    Returns documents with type='project_image'.
    """
    docs: list[Document] = []

    try:
        for image_path in sorted(projects_dir.glob("*.png")):
            if image_path.name in exclude_files:
                continue

            image_bytes = image_path.read_bytes()
            vision_caption = _try_image_vision_caption_bytes(
                image_bytes, filename=image_path.name
            )
            ocr_text = _try_image_ocr_bytes(image_bytes)
            filename = image_path.stem.replace("_", " ").replace("-", " ")
            caption = vision_caption or filename
            caption_label = "Vision caption" if vision_caption else "Filename"

            content = f"Image file: {image_path.name}. " f"{caption_label}: {caption}."
            if ocr_text:
                content += f" OCR text: {ocr_text}."

            doc = Document(
                page_content=content,
                metadata={
                    "source": str(image_path),
                    "type": "project_image",
                    "has_image": True,
                    "vision_caption": bool(vision_caption),
                    "caption_source": "vision" if vision_caption else "filename",
                    "ocr_text": bool(ocr_text),
                },
            )
            docs.append(doc)
            print(f"[projects] {image_path.name} -> image with caption")
    except Exception as e:
        print(f"[projects] Failed to load standalone images: {e}")

    return docs


# --------------------------------------------------------------------------- #
# Projects: Loading
# --------------------------------------------------------------------------- #


def load_projects(projects_dir: Path) -> list[Document]:
    """Load all project files (.md, .docx, .pptx, .png) from projects_dir.

    - Markdown files are split on H1/H2 headers for semantic chunking
    - DOCX files are extracted with heading hierarchy preserved as markdown headers
    - PPTX files are extracted with slide numbers as headers
    - Standalone PNG images are loaded with vision captions and OCR

    File stem becomes the `project` metadata key for filtering.
    """
    headers_to_split_on = [("#", "section"), ("##", "subsection")]
    md_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)

    docs: list[Document] = []
    if not projects_dir.exists():
        print(f"[projects] directory {projects_dir} does not exist")
        return docs

    # Track images referenced in markdown to avoid duplicates
    referenced_images: set[str] = set()

    # ========== Process Markdown files ==========
    for path in sorted(projects_dir.glob("*.md")):
        text = path.read_text(encoding="utf-8")
        if not text.strip():
            continue
        project_name = path.stem

        # Collect referenced images for later deduplication
        for alt, src, title in IMAGE_PATTERN.findall(text):
            if src and not src.startswith("http"):
                referenced_images.add(Path(src).name)

        image_docs = _collect_local_image_docs(text, path.parent, project_name)
        text = _normalize_markdown(text)

        chunks = md_splitter.split_text(text)
        for chunk in chunks:
            chunk.metadata.update(
                {
                    "source": str(path),
                    "project": project_name,
                    "type": "project",
                }
            )
            _chunk_metadata(chunk)
            docs.append(chunk)

        for image_doc in image_docs:
            _chunk_metadata(image_doc)
            docs.append(image_doc)

        print(
            f"[projects] {project_name}.md -> {len(chunks)} header-sections, {len(image_docs)} images"
        )

    # ========== Process DOCX files ==========
    for path in sorted(projects_dir.glob("*.docx")):
        text = _extract_docx_text(path)
        if not text.strip():
            continue
        project_name = path.stem

        text = _normalize_markdown(text)
        chunks = md_splitter.split_text(text)
        for chunk in chunks:
            chunk.metadata.update(
                {
                    "source": str(path),
                    "project": project_name,
                    "type": "project",
                }
            )
            _chunk_metadata(chunk)
            docs.append(chunk)

        print(f"[projects] {project_name}.docx -> {len(chunks)} header-sections")

    # ========== Process PPTX files ==========
    for path in sorted(projects_dir.glob("*.pptx")):
        text, pptx_image_bytes = _extract_pptx_text(path)
        if not text.strip() and not pptx_image_bytes:
            continue
        project_name = path.stem

        # Process slide text
        if text.strip():
            text = _normalize_markdown(text)
            chunks = md_splitter.split_text(text)
            for chunk in chunks:
                chunk.metadata.update(
                    {
                        "source": str(path),
                        "project": project_name,
                        "type": "project",
                    }
                )
                _chunk_metadata(chunk)
                docs.append(chunk)
        else:
            chunks = []

        # Process embedded images with vision captions
        for idx, image_bytes in enumerate(pptx_image_bytes):
            try:
                vision_caption = _try_image_vision_caption_bytes(
                    image_bytes, filename=f"{project_name}_slide_img_{idx}.png"
                )
                ocr_text = _try_image_ocr_bytes(image_bytes)
                caption = vision_caption or f"Image from {project_name}"
                caption_label = "Vision caption" if vision_caption else "Caption"

                content = (
                    f"Image from: {project_name}.pptx. " f"{caption_label}: {caption}."
                )
                if ocr_text:
                    content += f" OCR text: {ocr_text}."

                doc = Document(
                    page_content=content,
                    metadata={
                        "source": str(path),
                        "project": project_name,
                        "type": "project_image",
                        "has_image": True,
                        "vision_caption": bool(vision_caption),
                        "caption_source": "vision" if vision_caption else "pptx",
                        "ocr_text": bool(ocr_text),
                    },
                )
                _chunk_metadata(doc)
                docs.append(doc)
            except Exception as e:
                print(
                    f"[projects] Failed to process image {idx} from {project_name}.pptx: {e}"
                )

        print(
            f"[projects] {project_name}.pptx -> {len(chunks)} header-sections, {len(pptx_image_bytes)} images"
        )

    # ========== Process standalone PNG files ==========
    standalone_image_docs = _load_standalone_images(projects_dir, referenced_images)
    for image_doc in standalone_image_docs:
        _chunk_metadata(image_doc)
        docs.append(image_doc)

    if standalone_image_docs:
        print(f"[projects] loaded {len(standalone_image_docs)} standalone images")

    return docs


# --------------------------------------------------------------------------- #
# Chunking + indexing
# --------------------------------------------------------------------------- #
def chunk_documents(
    docs: list[Document],
    chunk_size: int = 800,
    overlap: int = 150,
) -> list[Document]:
    """Second-pass chunking — header sections can still be huge."""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return splitter.split_documents(docs)


def ingest_website(urls: list[str]) -> int:
    docs = scrape_website(urls)
    if not docs:
        print("[ingest] no website docs to index")
        return 0
    chunks = chunk_documents(docs)
    # Sanitize content to remove embedded instructions/code fences before indexing
    for c in chunks:
        try:
            c.page_content = sanitize_user_message(c.page_content)
        except Exception:
            pass
    vs = get_vectorstore(get_settings().website_collection)
    add_documents_safe(vs, chunks)
    print(f"[ingest] indexed {len(chunks)} website chunks from {len(urls)} URL(s)")
    return len(chunks)


def ingest_projects(projects_dir: Path) -> int:
    docs = load_projects(projects_dir)
    if not docs:
        print("[ingest] no project docs to index")
        return 0
    chunks = chunk_documents(docs)
    # Sanitize project chunks before indexing
    for c in chunks:
        try:
            c.page_content = sanitize_user_message(c.page_content)
        except Exception:
            pass
    vs = get_vectorstore(get_settings().projects_collection)
    add_documents_safe(vs, chunks)
    print(f"[ingest] indexed {len(chunks)} project chunks")
    return len(chunks)
